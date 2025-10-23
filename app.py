from flask import Flask, request, jsonify
from flask_cors import CORS
import sqlite3
import json
from datetime import datetime
import os
import requests
from werkzeug.utils import secure_filename
import PyPDF2
from docx import Document
from pptx import Presentation
import io
from ollama import Client

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# Ollama Configuration
OLLAMA_HOST = os.getenv('OLLAMA_HOST', 'https://ollama.com')
OLLAMA_API_KEY = os.getenv('OLLAMA_API_KEY', '1728cbe73f944db7afa1a3c8f52d2f41.GzEVZ8ADdcDHwIxdbvKnqbXy')
OLLAMA_MODEL = os.getenv('OLLAMA_MODEL', 'gpt-oss:120b-cloud')  # Cloud model

# Initialize Ollama client
ollama_client = Client(
    host=OLLAMA_HOST,
    headers={'Authorization': f'Bearer {OLLAMA_API_KEY}'}
)

# File Upload Configuration
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'doc', 'md', 'pptx', 'ppt'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200 MB max file size

# Database setup
# Database configuration
DATA_DIR = os.path.join(os.path.dirname(__file__), 'data')
DATABASE = os.path.join(DATA_DIR, 'nursing_app.db')

# Ensure data directory exists
os.makedirs(DATA_DIR, exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def get_db_connection():
    conn = sqlite3.connect(DATABASE)
    conn.row_factory = sqlite3.Row
    return conn

def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_stream):
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_docx(file_stream):
    """Extract text from DOCX file"""
    try:
        doc = Document(file_stream)
        text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def extract_text_from_pptx(file_stream):
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(file_stream)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        slide_text.append(row_text)

            text_content.append("\n".join(slide_text))

        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error extracting PowerPoint: {str(e)}"

def extract_text_from_file(file):
    """Extract text from uploaded file based on extension"""
    filename = secure_filename(file.filename)
    file_ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

    if file_ext == 'pdf':
        return extract_text_from_pdf(file)
    elif file_ext in ['docx', 'doc']:
        return extract_text_from_docx(file)
    elif file_ext in ['pptx', 'ppt']:
        return extract_text_from_pptx(file)
    elif file_ext in ['txt', 'md']:
        return file.read().decode('utf-8')
    else:
        return "Unsupported file type"

def init_database():
    """Initialize the database with required tables"""
    conn = get_db_connection()
    
    # Create tables
    conn.execute('''
        CREATE TABLE IF NOT EXISTS assignments (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            course TEXT NOT NULL,
            due_date TEXT NOT NULL,
            weight INTEGER,
            status TEXT DEFAULT 'not-started',
            completed BOOLEAN DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS clinical_shifts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            date TEXT NOT NULL,
            start_time TEXT,
            end_time TEXT,
            location TEXT NOT NULL,
            unit TEXT,
            hours REAL NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS requirements (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            deadline TEXT NOT NULL,
            status TEXT DEFAULT 'pending',
            renewal_months INTEGER DEFAULT 12,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS goals (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            description TEXT,
            target_date TEXT NOT NULL,
            category TEXT DEFAULT 'academic',
            completed BOOLEAN DEFAULT 0,
            created_date TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS grades (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            course TEXT NOT NULL,
            assessment TEXT NOT NULL,
            type TEXT DEFAULT 'assignment',
            grade REAL NOT NULL,
            max_points REAL NOT NULL,
            weight REAL,
            date TEXT,
            percentage REAL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS settings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            key TEXT UNIQUE NOT NULL,
            value TEXT NOT NULL,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS flashcards (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            question TEXT NOT NULL,
            answer TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.execute('''
        CREATE TABLE IF NOT EXISTS stress_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            date TEXT NOT NULL UNIQUE,
            stress_level INTEGER NOT NULL,
            mood TEXT,
            notes TEXT,
            study_hours REAL,
            sleep_hours REAL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    conn.execute('''
        CREATE TABLE IF NOT EXISTS weekly_activities (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            day_of_week TEXT NOT NULL,
            start_time TEXT NOT NULL,
            end_time TEXT NOT NULL,
            activity_type TEXT NOT NULL,
            title TEXT NOT NULL,
            description TEXT,
            week_start_date TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    conn.execute('''
        CREATE TABLE IF NOT EXISTS saved_tests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL,
            test_content TEXT NOT NULL,
            solutions_content TEXT NOT NULL,
            question_count INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    conn.execute('''
        CREATE TABLE IF NOT EXISTS test_attempts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            test_id INTEGER NOT NULL,
            mode TEXT NOT NULL CHECK(mode IN ('exam', 'practice')),
            score INTEGER,
            total_questions INTEGER NOT NULL,
            percentage REAL,
            time_spent_seconds INTEGER,
            timer_duration_minutes INTEGER,
            completed BOOLEAN DEFAULT 0,
            started_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            completed_at TIMESTAMP,
            FOREIGN KEY (test_id) REFERENCES saved_tests(id)
        )
    ''')

    conn.execute('''
        CREATE TABLE IF NOT EXISTS test_answers (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            attempt_id INTEGER NOT NULL,
            question_number INTEGER NOT NULL,
            user_answer TEXT,
            correct_answer TEXT,
            is_correct BOOLEAN,
            answered_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (attempt_id) REFERENCES test_attempts(id)
        )
    ''')

    conn.commit()
    conn.close()

# API Routes

@app.route('/api/assignments', methods=['GET', 'POST'])
def assignments():
    conn = get_db_connection()
    
    if request.method == 'GET':
        assignments = conn.execute('SELECT * FROM assignments ORDER BY due_date ASC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in assignments])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT INTO assignments (title, course, due_date, weight, status, completed)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (data['title'], data['course'], data['dueDate'],
              data.get('weight'), data.get('status', 'not-started'),
              data.get('completed', False)))
        conn.commit()
        assignment_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': assignment_id}), 201

@app.route('/api/assignments/<int:assignment_id>', methods=['PUT', 'DELETE'])
def assignment_detail(assignment_id):
    conn = get_db_connection()
    
    if request.method == 'PUT':
        data = request.json
        conn.execute('''
            UPDATE assignments 
            SET title=?, course=?, due_date=?, weight=?, status=?, completed=?
            WHERE id=?
        ''', (data['title'], data['course'], data['dueDate'], 
              data.get('weight'), data.get('status'), 
              data.get('completed'), assignment_id))
        conn.commit()
        conn.close()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        conn.execute('DELETE FROM assignments WHERE id = ?', (assignment_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True})

@app.route('/api/clinical-shifts', methods=['GET', 'POST'])
def clinical_shifts():
    conn = get_db_connection()
    
    if request.method == 'GET':
        shifts = conn.execute('SELECT * FROM clinical_shifts ORDER BY date ASC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in shifts])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT INTO clinical_shifts (date, start_time, end_time, location, unit, hours)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (data['date'], data.get('startTime'), data.get('endTime'),
              data['location'], data.get('unit'), data['hours']))
        conn.commit()
        shift_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': shift_id}), 201

@app.route('/api/clinical-shifts/<int:shift_id>', methods=['PUT', 'DELETE'])
def clinical_shift_detail(shift_id):
    conn = get_db_connection()
    
    if request.method == 'PUT':
        data = request.json
        conn.execute('''
            UPDATE clinical_shifts 
            SET date=?, start_time=?, end_time=?, location=?, unit=?, hours=?
            WHERE id=?
        ''', (data['date'], data.get('start_time'), data.get('end_time'),
              data['location'], data.get('unit'), data['hours'], shift_id))
        conn.commit()
        conn.close()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        conn.execute('DELETE FROM clinical_shifts WHERE id = ?', (shift_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True})

@app.route('/api/requirements', methods=['GET', 'POST'])
def requirements():
    conn = get_db_connection()
    
    if request.method == 'GET':
        requirements = conn.execute('SELECT * FROM requirements ORDER BY deadline ASC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in requirements])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT INTO requirements (name, deadline, status, renewal_months)
            VALUES (?, ?, ?, ?)
        ''', (data['name'], data['deadline'], data.get('status', 'pending'),
              data.get('renewalMonths', 12)))
        conn.commit()
        req_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': req_id}), 201

@app.route('/api/requirements/<int:req_id>', methods=['PUT', 'DELETE'])
def requirement_detail(req_id):
    conn = get_db_connection()
    
    if request.method == 'PUT':
        data = request.json
        conn.execute('''
            UPDATE requirements 
            SET name=?, deadline=?, status=?, renewal_months=?
            WHERE id=?
        ''', (data['name'], data['deadline'], data['status'],
              data.get('renewalMonths', 12), req_id))
        conn.commit()
        conn.close()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        conn.execute('DELETE FROM requirements WHERE id = ?', (req_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True})

@app.route('/api/goals', methods=['GET', 'POST'])
def goals():
    conn = get_db_connection()
    
    if request.method == 'GET':
        goals = conn.execute('SELECT * FROM goals ORDER BY target_date ASC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in goals])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT INTO goals (title, description, target_date, category, completed, created_date)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (data['title'], data.get('description'), data['targetDate'],
              data.get('category', 'academic'), data.get('completed', False),
              data.get('createdDate')))
        conn.commit()
        goal_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': goal_id}), 201

@app.route('/api/goals/<int:goal_id>', methods=['PUT', 'DELETE'])
def goal_detail(goal_id):
    conn = get_db_connection()
    
    if request.method == 'PUT':
        data = request.json
        conn.execute('''
            UPDATE goals 
            SET title=?, description=?, target_date=?, category=?, completed=?
            WHERE id=?
        ''', (data['title'], data.get('description'), data['targetDate'],
              data.get('category'), data.get('completed'), goal_id))
        conn.commit()
        conn.close()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        conn.execute('DELETE FROM goals WHERE id = ?', (goal_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True})

@app.route('/api/grades', methods=['GET', 'POST'])
def grades():
    conn = get_db_connection()
    
    if request.method == 'GET':
        grades = conn.execute('SELECT * FROM grades ORDER BY date DESC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in grades])
    
    elif request.method == 'POST':
        data = request.json
        percentage = round((float(data['grade']) / float(data['maxPoints'])) * 100)
        cursor = conn.execute('''
            INSERT INTO grades (course, assessment, type, grade, max_points, weight, date, percentage)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        ''', (data['course'], data['assessment'], data.get('type', 'assignment'),
              data['grade'], data['maxPoints'], data.get('weight'),
              data.get('date'), percentage))
        conn.commit()
        grade_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': grade_id}), 201

@app.route('/api/grades/<int:grade_id>', methods=['DELETE'])
def grade_detail(grade_id):
    conn = get_db_connection()
    conn.execute('DELETE FROM grades WHERE id = ?', (grade_id,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/api/flashcards', methods=['GET', 'POST'])
def flashcards():
    conn = get_db_connection()
    
    if request.method == 'GET':
        flashcards = conn.execute('SELECT * FROM flashcards ORDER BY created_at DESC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in flashcards])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT INTO flashcards (question, answer)
            VALUES (?, ?)
        ''', (data['question'], data['answer']))
        conn.commit()
        card_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': card_id}), 201

@app.route('/api/flashcards/<int:card_id>', methods=['DELETE'])
def flashcard_detail(card_id):
    conn = get_db_connection()
    conn.execute('DELETE FROM flashcards WHERE id = ?', (card_id,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/api/stress-logs', methods=['GET', 'POST'])
def stress_logs():
    conn = get_db_connection()
    
    if request.method == 'GET':
        logs = conn.execute('SELECT * FROM stress_logs ORDER BY date DESC').fetchall()
        conn.close()
        return jsonify([dict(row) for row in logs])
    
    elif request.method == 'POST':
        data = request.json
        cursor = conn.execute('''
            INSERT OR REPLACE INTO stress_logs (date, stress_level, mood, notes, study_hours, sleep_hours)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (data['date'], data['stress_level'], data.get('mood'),
              data.get('notes'), data.get('study_hours'), data.get('sleep_hours')))
        conn.commit()
        log_id = cursor.lastrowid
        conn.close()
        return jsonify({'id': log_id}), 201

@app.route('/api/stress-logs/<int:log_id>', methods=['DELETE'])
def stress_log_detail(log_id):
    conn = get_db_connection()
    conn.execute('DELETE FROM stress_logs WHERE id = ?', (log_id,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/api/settings', methods=['GET', 'POST'])
def settings():
    conn = get_db_connection()
    
    if request.method == 'GET':
        settings = conn.execute('SELECT * FROM settings').fetchall()
        conn.close()
        return jsonify({row['key']: row['value'] for row in settings})
    
    elif request.method == 'POST':
        data = request.json
        for key, value in data.items():
            conn.execute('''
                INSERT OR REPLACE INTO settings (key, value, updated_at)
                VALUES (?, ?, CURRENT_TIMESTAMP)
            ''', (key, str(value)))
        conn.commit()
        conn.close()
        return jsonify({'success': True})

@app.route('/api/initialize', methods=['POST'])
def initialize_data():
    """Initialize database with default requirements"""
    conn = get_db_connection()
    
    # Check if requirements already exist
    existing = conn.execute('SELECT COUNT(*) as count FROM requirements').fetchone()
    
    if existing['count'] == 0:
        # Add default requirements
        default_requirements = [
            ('ERV Letter (Year 1)', '2025-08-01', 'pending', 12),
            ('ERV Letter (Year 2)', '2026-08-01', 'pending', 12),
            ('Health Assessment Record', '2025-08-01', 'pending', 12),
            ('CPR-BLS', '2025-08-01', 'pending', 12),
            ('Standard First Aid', '2025-08-01', 'pending', 36),
            ('Mask Fit', '2025-08-01', 'pending', 24),
            ('Police Vulnerable Sector Check', '2025-08-01', 'pending', 12)
        ]
        
        for req in default_requirements:
            conn.execute('''
                INSERT INTO requirements (name, deadline, status, renewal_months)
                VALUES (?, ?, ?, ?)
            ''', req)
        
        # Set default current semester
        conn.execute('''
            INSERT OR REPLACE INTO settings (key, value)
            VALUES ('currentSemester', '1')
        ''')
        
        conn.commit()
    
    conn.close()
    return jsonify({'success': True})

@app.route('/')
def index():
    """Serve the main application HTML file"""
    with open('database_enabled_frontend.html', 'r', encoding='utf-8') as f:
        return f.read()

@app.route('/favicon_io/<path:filename>')
def favicon_files(filename):
    """Serve favicon files"""
    from flask import send_from_directory
    return send_from_directory('favicon_io', filename)

@app.route('/favicon.ico')
def favicon():
    """Serve favicon.ico from root"""
    from flask import send_from_directory
    return send_from_directory('favicon_io', 'favicon.ico')

@app.route('/api/backup', methods=['GET'])
def backup_data():
    """Export all data as JSON for backup"""
    conn = get_db_connection()
    
    backup = {
        'assignments': [dict(row) for row in conn.execute('SELECT * FROM assignments').fetchall()],
        'clinical_shifts': [dict(row) for row in conn.execute('SELECT * FROM clinical_shifts').fetchall()],
        'requirements': [dict(row) for row in conn.execute('SELECT * FROM requirements').fetchall()],
        'goals': [dict(row) for row in conn.execute('SELECT * FROM goals').fetchall()],
        'grades': [dict(row) for row in conn.execute('SELECT * FROM grades').fetchall()],
        'flashcards': [dict(row) for row in conn.execute('SELECT * FROM flashcards').fetchall()],
        'settings': {row['key']: row['value'] for row in conn.execute('SELECT * FROM settings').fetchall()},
        'backup_date': datetime.now().isoformat()
    }
    
    conn.close()
    return jsonify(backup)

@app.route('/api/generate-test', methods=['POST'])
def generate_test():
    """Generate practice test and solution sheet from uploaded study materials using Ollama AI"""
    try:
        print("=== Starting test generation ===")

        # Check if files were uploaded
        if 'files' not in request.files:
            print("ERROR: No files in request")
            return jsonify({'error': 'No files uploaded'}), 400

        files = request.files.getlist('files')
        prompt = request.form.get('prompt', 'Generate a comprehensive practice test')
        print(f"Received {len(files)} files")

        if not files or all(file.filename == '' for file in files):
            print("ERROR: No files selected or empty filenames")
            return jsonify({'error': 'No files selected'}), 400

        # Extract text from all uploaded files
        study_materials = []
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                print(f"Extracting text from: {filename}")
                text = extract_text_from_file(file)
                print(f"Extracted {len(text)} characters from {filename}")
                study_materials.append({
                    'filename': filename,
                    'content': text
                })

        if not study_materials:
            print("ERROR: No valid files after processing")
            return jsonify({'error': 'No valid files uploaded'}), 400

        print(f"Successfully processed {len(study_materials)} files")

        # Combine all study materials
        combined_content = "\n\n---\n\n".join([
            f"File: {mat['filename']}\n\n{mat['content']}"
            for mat in study_materials
        ])
        print(f"Total content length: {len(combined_content)} characters")

        # Build the prompt for Ollama
        system_prompt = f"""You are a nursing education assistant. Based on the provided study materials, create a comprehensive practice test and separate solution sheet.

IMPORTANT USER REQUEST: {prompt}

Study Materials:
{combined_content}

Generate TWO separate documents in markdown format following the user's request exactly:

1. TEST DOCUMENT:
   - Follow the user's instructions for the number and type of questions
   - If user specified a number (e.g., "100 questions"), create that many questions
   - Mix question types: Multiple choice, True/False, Short answer, Fill-in-the-blank, Matching
   - Make questions challenging but fair
   - Cover ALL chapters/topics mentioned in the materials
   - Focus on key concepts, terminology, processes, and clinical applications

2. SOLUTION SHEET:
   - Provide answers to ALL questions
   - Include brief explanations for each answer
   - Reference specific slides/sections when possible

Format your response as JSON:
{{
  "test": "# Practice Test\\n\\n## Multiple Choice (Questions 1-X)\\n1. Question text...\\n\\n## True/False (Questions X-Y)\\n...",
  "solutions": "# Solution Sheet\\n\\n## Multiple Choice Answers\\n1. Answer: [Correct answer] - Explanation...\\n\\n## True/False Answers\\n..."
}}

Respond ONLY with valid JSON. Make sure to honor the user's requested number of questions."""

        # Call Ollama Cloud API
        print(f"Calling Ollama Cloud at: {OLLAMA_HOST}")
        print(f"Using model: {OLLAMA_MODEL}")
        print(f"Prompt length: {len(system_prompt)} characters")

        # Try with retry logic for cloud API
        max_retries = 2
        for attempt in range(max_retries):
            try:
                print(f"Attempt {attempt + 1} of {max_retries}")

                # Use ollama client library to generate response
                response = ollama_client.generate(
                    model=OLLAMA_MODEL,
                    prompt=system_prompt,
                    options={
                        'temperature': 0.7,
                        'num_predict': 16000  # Large limit for 100+ questions (~16k tokens)
                    }
                )

                generated_text = response['response']
                print(f"✓ Received response: {len(generated_text)} characters")
                break  # Success, exit retry loop

            except Exception as e:
                error_msg = str(e)
                print(f"✗ Attempt {attempt + 1} failed: {error_msg}")

                if attempt < max_retries - 1:
                    print(f"Retrying in 3 seconds...")
                    import time
                    time.sleep(3)
                else:
                    # Last attempt failed
                    print(f"All {max_retries} attempts failed")

                    # Provide helpful error message
                    if '502' in error_msg or 'upstream' in error_msg:
                        return jsonify({
                            'error': 'Ollama Cloud is temporarily overloaded',
                            'details': 'The cloud service returned a 502 error. This usually happens with large requests. Try: 1) Uploading fewer files, 2) Requesting fewer questions (e.g., 50 instead of 100), or 3) Waiting a few minutes and trying again.',
                            'technical_details': error_msg
                        }), 500
                    else:
                        return jsonify({
                            'error': 'Failed to connect to Ollama Cloud',
                            'details': str(e),
                            'host': OLLAMA_HOST
                        }), 500

        # Try to parse as JSON first
        try:
            # Look for JSON in the response
            json_start = generated_text.find('{')
            json_end = generated_text.rfind('}') + 1
            if json_start != -1 and json_end > json_start:
                json_text = generated_text[json_start:json_end]
                test_data = json.loads(json_text)
                return jsonify(test_data)
            else:
                # If no JSON found, split the response manually
                parts = generated_text.split('SOLUTION SHEET')
                test_content = parts[0].replace('TEST DOCUMENT', '').strip()
                solution_content = parts[1].strip() if len(parts) > 1 else "Solutions not generated properly."

                return jsonify({
                    'test': test_content,
                    'solutions': solution_content
                })
        except (json.JSONDecodeError, IndexError):
            # If all parsing fails, return raw response
            return jsonify({
                'test': generated_text,
                'solutions': 'Please review the test document above for answers.',
                'note': 'AI response format was not as expected'
            })

    except requests.exceptions.RequestException as e:
        return jsonify({
            'error': 'Failed to connect to Ollama',
            'details': str(e),
            'hint': 'Make sure Ollama is running with: ollama serve'
        }), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Saved Tests API Routes

@app.route('/api/tests/save', methods=['POST'])
def save_test():
    """Save a generated test to the database"""
    try:
        data = request.json
        print(f"Received save request with data keys: {data.keys()}")
        print(f"Title: {data.get('title', 'MISSING')}")
        print(f"Test length: {len(data.get('test', ''))}")
        print(f"Solutions length: {len(data.get('solutions', ''))}")

        conn = get_db_connection()

        cursor = conn.execute('''
            INSERT INTO saved_tests (title, test_content, solutions_content, question_count)
            VALUES (?, ?, ?, ?)
        ''', (data['title'], data['test'], data['solutions'], data.get('questionCount', 0)))

        conn.commit()
        test_id = cursor.lastrowid
        conn.close()

        print(f"Successfully saved test with ID: {test_id}")
        return jsonify({'id': test_id, 'success': True}), 201
    except Exception as e:
        print(f"ERROR in save_test: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests', methods=['GET'])
def get_tests():
    """Get all saved tests"""
    try:
        conn = get_db_connection()
        tests = conn.execute('''
            SELECT st.*,
                   COUNT(DISTINCT ta.id) as attempt_count,
                   MAX(ta.percentage) as best_score
            FROM saved_tests st
            LEFT JOIN test_attempts ta ON st.id = ta.test_id AND ta.completed = 1
            GROUP BY st.id
            ORDER BY st.created_at DESC
        ''').fetchall()
        conn.close()

        return jsonify([dict(row) for row in tests])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/<int:test_id>', methods=['GET', 'DELETE'])
def test_detail(test_id):
    """Get or delete a specific test"""
    conn = get_db_connection()

    if request.method == 'GET':
        try:
            test = conn.execute('SELECT * FROM saved_tests WHERE id = ?', (test_id,)).fetchone()
            conn.close()

            if test:
                return jsonify(dict(test))
            else:
                return jsonify({'error': 'Test not found'}), 404
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    elif request.method == 'DELETE':
        try:
            # Delete test and all related attempts/answers (cascade)
            conn.execute('DELETE FROM test_answers WHERE attempt_id IN (SELECT id FROM test_attempts WHERE test_id = ?)', (test_id,))
            conn.execute('DELETE FROM test_attempts WHERE test_id = ?', (test_id,))
            conn.execute('DELETE FROM saved_tests WHERE id = ?', (test_id,))
            conn.commit()
            conn.close()

            return jsonify({'success': True})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

# Test Attempts API Routes

@app.route('/api/tests/<int:test_id>/start', methods=['POST'])
def start_test_attempt(test_id):
    """Start a new test attempt"""
    try:
        data = request.json
        conn = get_db_connection()

        # Verify test exists
        test = conn.execute('SELECT question_count FROM saved_tests WHERE id = ?', (test_id,)).fetchone()
        if not test:
            conn.close()
            return jsonify({'error': 'Test not found'}), 404

        cursor = conn.execute('''
            INSERT INTO test_attempts (test_id, mode, total_questions, timer_duration_minutes)
            VALUES (?, ?, ?, ?)
        ''', (test_id, data['mode'], test['question_count'], data.get('timerDuration')))

        conn.commit()
        attempt_id = cursor.lastrowid
        conn.close()

        return jsonify({'id': attempt_id, 'success': True}), 201
    except Exception as e:
        print(f"ERROR in start_test_attempt: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/attempts/<int:attempt_id>', methods=['GET', 'PUT', 'DELETE'])
def test_attempt_detail(attempt_id):
    """Get, update, or delete a test attempt (for practice mode resume)"""
    conn = get_db_connection()

    if request.method == 'GET':
        try:
            attempt = conn.execute('SELECT * FROM test_attempts WHERE id = ?', (attempt_id,)).fetchone()
            if not attempt:
                conn.close()
                return jsonify({'error': 'Attempt not found'}), 404

            # Get saved answers
            answers = conn.execute('SELECT * FROM test_answers WHERE attempt_id = ? ORDER BY question_number', (attempt_id,)).fetchall()
            conn.close()

            return jsonify({
                'attempt': dict(attempt),
                'answers': [dict(row) for row in answers]
            })
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    elif request.method == 'PUT':
        try:
            data = request.json

            # Save/update answer
            if 'answer' in data:
                conn.execute('''
                    INSERT OR REPLACE INTO test_answers (attempt_id, question_number, user_answer, correct_answer, is_correct)
                    VALUES (?, ?, ?, ?, ?)
                ''', (attempt_id, data['questionNumber'], data['answer'], data.get('correctAnswer'), data.get('isCorrect')))

            conn.commit()
            conn.close()

            return jsonify({'success': True})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    elif request.method == 'DELETE':
        try:
            # Delete all answers for this attempt
            conn.execute('DELETE FROM test_answers WHERE attempt_id = ?', (attempt_id,))
            # Delete the attempt
            conn.execute('DELETE FROM test_attempts WHERE id = ?', (attempt_id,))
            conn.commit()
            conn.close()

            return jsonify({'success': True})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

@app.route('/api/tests/<int:test_id>/attempts', methods=['GET'])
def get_test_attempts(test_id):
    """Get all attempts for a test"""
    try:
        conn = get_db_connection()
        attempts = conn.execute('''
            SELECT * FROM test_attempts
            WHERE test_id = ?
            ORDER BY started_at DESC
        ''', (test_id,)).fetchall()
        conn.close()

        return jsonify([dict(row) for row in attempts])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/attempts/<int:attempt_id>/answers', methods=['GET'])
def get_attempt_answers(attempt_id):
    """Get all saved answers for an attempt"""
    try:
        conn = get_db_connection()
        answers = conn.execute('''
            SELECT * FROM test_answers
            WHERE attempt_id = ?
            ORDER BY question_number ASC
        ''', (attempt_id,)).fetchall()
        conn.close()

        return jsonify([dict(row) for row in answers])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/attempts/<int:attempt_id>/submit', methods=['POST'])
def submit_test_attempt(attempt_id):
    """Submit a completed test attempt"""
    try:
        data = request.json
        conn = get_db_connection()

        # Calculate score
        score = data['score']
        total = data['total']
        percentage = round((score / total) * 100) if total > 0 else 0

        # Update attempt
        conn.execute('''
            UPDATE test_attempts
            SET score = ?, percentage = ?, time_spent_seconds = ?, completed = 1, completed_at = CURRENT_TIMESTAMP
            WHERE id = ?
        ''', (score, percentage, data.get('timeSpent', 0), attempt_id))

        # Save all answers if not already saved (for exam mode)
        if 'answers' in data:
            for answer in data['answers']:
                conn.execute('''
                    INSERT OR REPLACE INTO test_answers (attempt_id, question_number, user_answer, correct_answer, is_correct)
                    VALUES (?, ?, ?, ?, ?)
                ''', (attempt_id, answer['questionNumber'], answer['userAnswer'],
                      answer['correctAnswer'], answer['isCorrect']))

        conn.commit()
        conn.close()

        return jsonify({'success': True, 'score': score, 'percentage': percentage})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/<int:test_id>/history', methods=['GET'])
def get_test_history(test_id):
    """Get all attempts and scores for a test"""
    try:
        conn = get_db_connection()
        attempts = conn.execute('''
            SELECT * FROM test_attempts
            WHERE test_id = ? AND completed = 1
            ORDER BY completed_at DESC
        ''', (test_id,)).fetchall()
        conn.close()

        return jsonify([dict(row) for row in attempts])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/tests/<int:test_id>/analytics', methods=['GET'])
def get_test_analytics(test_id):
    """Get analytics data for a test"""
    try:
        conn = get_db_connection()

        # Get all completed attempts
        attempts = conn.execute('''
            SELECT * FROM test_attempts
            WHERE test_id = ? AND completed = 1
            ORDER BY completed_at ASC
        ''', (test_id,)).fetchall()

        if not attempts:
            conn.close()
            return jsonify({'error': 'No completed attempts found'}), 404

        # Calculate statistics
        scores = [a['percentage'] for a in attempts]
        best_score = max(scores)
        avg_score = sum(scores) / len(scores)
        latest_score = scores[-1]

        # Get performance by question
        question_performance = conn.execute('''
            SELECT question_number,
                   SUM(CASE WHEN is_correct = 1 THEN 1 ELSE 0 END) as correct_count,
                   COUNT(*) as total_attempts,
                   ROUND(SUM(CASE WHEN is_correct = 1 THEN 1 ELSE 0 END) * 100.0 / COUNT(*), 2) as accuracy
            FROM test_answers
            WHERE attempt_id IN (SELECT id FROM test_attempts WHERE test_id = ? AND completed = 1)
            GROUP BY question_number
            ORDER BY accuracy ASC
        ''', (test_id,)).fetchall()

        conn.close()

        return jsonify({
            'attempts': [dict(row) for row in attempts],
            'statistics': {
                'bestScore': best_score,
                'averageScore': round(avg_score, 1),
                'latestScore': latest_score,
                'totalAttempts': len(attempts)
            },
            'questionPerformance': [dict(row) for row in question_performance]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # Initialize database on startup
    init_database()
    
    # Run the app
    app.run(debug=True, host='0.0.0.0', port=5008)